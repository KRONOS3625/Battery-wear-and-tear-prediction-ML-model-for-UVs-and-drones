from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
MODELS = ROOT / "models"
OUTPUT = ROOT / "Battery_Intelligence_Project_Presentation.pptx"


BG = RGBColor(6, 13, 28)
PANEL = RGBColor(17, 28, 49)
ACCENT = RGBColor(46, 213, 191)
ACCENT_2 = RGBColor(249, 115, 22)
TEXT = RGBColor(245, 248, 255)
MUTED = RGBColor(189, 204, 225)
SOFT = RGBColor(96, 119, 155)


with (MODELS / "model_metadata.json").open("r", encoding="utf-8") as handle:
    metadata = json.load(handle)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.7), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.1), Inches(11.4), Inches(0.4))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(11)
        run2.font.color.rgb = MUTED


def add_panel(slide, x: float, y: float, w: float, h: float, title: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(44, 61, 96)
    shape.line.width = Pt(1.2)
    if title:
        tx = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.08), Inches(w - 0.4), Inches(0.3))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = ACCENT
    return shape


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, font_size: int = 16, color=TEXT) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_two_col_table_like(slide, rows: list[tuple[str, str]], x: float, y: float, w: float, row_h: float = 0.38) -> None:
    for idx, (left, right) in enumerate(rows):
        ypos = y + idx * row_h
        left_box = slide.shapes.add_textbox(Inches(x), Inches(ypos), Inches(w * 0.38), Inches(row_h))
        right_box = slide.shapes.add_textbox(Inches(x + w * 0.4), Inches(ypos), Inches(w * 0.58), Inches(row_h))
        for box, text, col, bold in [(left_box, left, MUTED, True), (right_box, right, TEXT, False)]:
            tf = box.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.size = Pt(13)
            r.font.color.rgb = col
            r.font.bold = bold


def add_metric_card(slide, label: str, value: str, detail: str, x: float, y: float, w: float, h: float) -> None:
    add_panel(slide, x, y, w, h)
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.16), Inches(w - 0.35), Inches(h - 0.2))
    tf = box.text_frame
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = label
    r1.font.size = Pt(11)
    r1.font.color.rgb = SOFT
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = TEXT
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = detail
    r3.font.size = Pt(10)
    r3.font.color.rgb = MUTED


def add_image(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.3))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    add_title(slide, "Battery Intelligence Dashboard for EVs and Drones", "Machine-learning based prediction of battery wear, degradation, and remaining useful life")
    add_bullets(
        slide,
        [
            "Uses real NASA Li-ion battery aging data stored in the project workspace",
            "Combines ML prediction, anomaly detection, and realistic EV/drone operating stress simulation",
            "Predicts SoH, wear %, RUL, end-of-life, failure probability, thermal risk, and pack imbalance",
        ],
        0.8,
        1.8,
        6.2,
        2.4,
        font_size=18,
    )
    add_metric_card(slide, "Training Samples", str(metadata["sample_count"]), "Real discharge records", 8.0, 1.8, 2.0, 1.4)
    add_metric_card(slide, "Battery Cells", str(metadata["battery_count"]), "Unique NASA cells", 10.3, 1.8, 2.0, 1.4)
    add_metric_card(slide, "SoH MAE", f'{metadata["metrics"]["soh_mae"]:.2f}%', "Prediction error", 8.0, 3.45, 2.0, 1.4)
    add_metric_card(slide, "RUL MAE", f'{metadata["metrics"]["rul_mae"]:.2f}', "Cycles error", 10.3, 3.45, 2.0, 1.4)
    add_bullets(slide, ["Team project presentation deck generated automatically from the current codebase"], 0.8, 6.8, 11.0, 0.5, font_size=12, color=MUTED)


def abbreviations_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Shortforms Used", "Abbreviations used across the model, UI, formulas, and explanation")
    add_panel(slide, 0.6, 1.45, 12.1, 5.6)
    rows = [
        ("EV", "Electric Vehicle"),
        ("SoH", "State of Health"),
        ("RUL", "Remaining Useful Life"),
        ("DoD", "Depth of Discharge"),
        ("SoC", "State of Charge"),
        ("MAE", "Mean Absolute Error"),
        ("R²", "Coefficient of Determination"),
        ("C-rate", "Charge/Discharge rate normalized by battery capacity"),
        ("Wh", "Watt-hour, energy delivered by the battery"),
        ("ML", "Machine Learning"),
        ("UI", "User Interface"),
        ("Li-ion / LiPo / LFP", "Battery chemistry families used in the scenario layer"),
    ]
    add_two_col_table_like(slide, rows, 1.0, 1.85, 10.8, row_h=0.42)


def workflow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "How The Project Works", "End-to-end workflow from raw NASA battery files to dashboard output")
    steps = [
        ("1. Dataset ingestion", "MATLAB .mat files are read from NASA battery aging archives in the workspace."),
        ("2. Feature engineering", "Discharge capacity, voltage, current, temperature, energy, impedance, and C-rate are extracted."),
        ("3. Label generation", "SoH and RUL labels are derived from each battery’s degradation trajectory."),
        ("4. Model comparison", "Random Forest, Extra Trees, and Gradient Boosting are trained and benchmarked."),
        ("5. Final prediction layer", "Best models plus anomaly detection are saved and loaded by the web server."),
        ("6. Scenario realism layer", "EV/drone profile, chemistry, SoC, DoD, mission stress, age, and fast charging adjust the output."),
        ("7. Dashboard output", "UI shows SoH, wear, RUL, failure risk, heatmap, contributions, and what-if recommendations."),
    ]
    for idx, (label, desc) in enumerate(steps):
        y = 1.5 + idx * 0.76
        add_panel(slide, 0.7, y, 2.2, 0.58)
        add_panel(slide, 3.1, y, 9.3, 0.58)
        add_bullets(slide, [label], 0.92, y + 0.12, 1.8, 0.3, font_size=14)
        add_bullets(slide, [desc], 3.3, y + 0.1, 8.9, 0.35, font_size=14, color=MUTED)


def data_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Dataset And Training Base", "Real battery-aging data used to build the ML core")
    add_panel(slide, 0.7, 1.45, 5.7, 4.8, "Dataset Facts")
    add_bullets(
        slide,
        [
            f'Training samples: {metadata["sample_count"]}',
            f'Battery cells: {metadata["battery_count"]}',
            "Source: NASA Li-ion Battery Aging Dataset and randomized battery usage series",
            "Local workspace contains zipped MATLAB datasets and readme files",
            "Rows are built from discharge and impedance cycles only",
            "Invalid or nonphysical rows are filtered before training",
        ],
        0.95,
        1.95,
        4.9,
        3.8,
        font_size=16,
    )
    add_panel(slide, 6.7, 1.45, 5.9, 4.8, "Learned Feature Set")
    add_bullets(
        slide,
        [
            "Internal resistance",
            "Capacity",
            "Cycle number",
            "Cell and ambient temperature",
            "Average and minimum voltage",
            "Average / peak current and current variation",
            "Discharge duration, energy delivered, and load C-rate",
        ],
        6.95,
        1.95,
        5.2,
        3.9,
        font_size=16,
    )


def parameters_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Parameters Used And Their Effect", "How each major input affects the output prediction")
    rows = [
        ("Internal resistance", "Higher resistance usually increases wear, lowers SoH, and raises failure risk."),
        ("Capacity", "Lower available capacity directly reduces SoH and shortens RUL."),
        ("Cycle number", "More cycles generally mean older battery state and lower remaining life."),
        ("Temperature", "High temperature accelerates degradation and pushes thermal risk upward."),
        ("Ambient temperature", "Extreme ambient conditions increase stress even if cell temperature looks safe."),
        ("Voltage under load", "More voltage sag suggests internal stress and possible pack weakness."),
        ("Current / C-rate", "Higher current increases thermal load and charge/discharge severity."),
        ("SoC", "Very high or very low SoC can increase stress and aging."),
        ("DoD", "Deep discharge cycles generally increase long-term degradation."),
        ("Age / fast charging / mission stress", "Calendar aging and harsh usage lower SoH and increase failure probability."),
    ]
    add_panel(slide, 0.7, 1.45, 12.0, 5.7)
    add_two_col_table_like(slide, rows, 1.0, 1.8, 11.2, row_h=0.5)


def formulas_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Core Formulas Used", "Main equations used in feature engineering, labeling, and output generation")
    left = [
        "1. Load C-rate = Average Current / Capacity",
        "Why: normalizes current by battery size and measures how aggressively the battery is being used.",
        "",
        "2. Energy (Wh) = Integral[ Voltage x |Current| dt ] / 3600",
        "Why: captures actual energy delivered during discharge, not just instantaneous values.",
        "",
        "3. SoH (%) = (Current Capacity / Initial or Best Capacity) x 100",
        "Why: standard battery-health definition used to quantify degradation.",
        "",
        "4. RUL (cycles) = End-of-Life Cycle - Current Cycle",
        "Why: directly measures remaining life before the battery crosses the end-of-life threshold.",
    ]
    right = [
        "5. Voltage Sag Risk = (Average Voltage - Minimum Voltage) / 0.7",
        "Why: larger sag often indicates internal resistance rise or stressed cells.",
        "",
        "6. Failure Probability = Sigmoid( wear + stress + anomaly + temperature effects )",
        "Why: converts several risk signals into an interpretable probability-like output.",
        "",
        "7. Degradation Rate per Cycle = Wear % / Cycle Number",
        "Why: gives a quick estimate of how rapidly the pack is aging.",
        "",
        "8. Stress Index = weighted sum of SoC, DoD, age, fast-charge, vibration, distance, speed, and thermal factors",
        "Why: adds real-world EV/drone behavior that does not fully exist in the raw training dataset.",
    ]
    add_panel(slide, 0.7, 1.5, 5.9, 5.6)
    add_panel(slide, 6.75, 1.5, 5.9, 5.6)
    add_bullets(slide, left, 0.95, 1.8, 5.2, 5.0, font_size=14)
    add_bullets(slide, right, 6.98, 1.8, 5.2, 5.0, font_size=14)


def algorithms_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "ML Algorithms Used And Why", "All prediction and detection algorithms used in the project")
    algo_rows = [
        ("RandomForestRegressor", "Strong ensemble baseline for nonlinear tabular battery data. Robust and easy to interpret."),
        ("ExtraTreesRegressor", "Selected final regressor. Gives strong performance on mixed engineered battery features."),
        ("GradientBoostingRegressor", "Useful comparison model that sequentially improves prediction residuals."),
        ("IsolationForest", "Detects abnormal battery states that look unusual compared with training data."),
    ]
    add_panel(slide, 0.7, 1.45, 6.1, 5.7, "Algorithms")
    add_two_col_table_like(slide, algo_rows, 0.95, 1.9, 5.45, row_h=0.7)
    add_panel(slide, 7.0, 1.45, 5.6, 5.7, "Why Model Comparison Matters")
    add_bullets(
        slide,
        [
            "Different algorithms behave differently on battery-health tabular data.",
            "The training script evaluates MAE and R² for SoH and RUL separately.",
            f'Selected SoH model: {metadata["model_comparison"]["soh"]["selected"]["name"]}',
            f'Selected RUL model: {metadata["model_comparison"]["rul"]["selected"]["name"]}',
            "This makes the system more defensible than choosing a single model without comparison.",
        ],
        7.25,
        1.9,
        4.9,
        3.9,
        font_size=15,
    )


def results_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Training Results", "Current performance from the latest run in this workspace")
    add_metric_card(slide, "SoH MAE", f'{metadata["metrics"]["soh_mae"]:.2f}%', "Lower is better", 0.8, 1.5, 2.4, 1.35)
    add_metric_card(slide, "SoH R²", f'{metadata["metrics"]["soh_r2"]:.3f}', "Closer to 1 is better", 3.4, 1.5, 2.4, 1.35)
    add_metric_card(slide, "RUL MAE", f'{metadata["metrics"]["rul_mae"]:.2f}', "Cycles error", 6.0, 1.5, 2.4, 1.35)
    add_metric_card(slide, "RUL R²", f'{metadata["metrics"]["rul_r2"]:.3f}', "Closer to 1 is better", 8.6, 1.5, 2.4, 1.35)
    add_image(slide, ASSETS / "model_diagnostics.png", 0.8, 3.0, 5.8, 3.5)
    add_image(slide, ASSETS / "feature_importance.png", 6.9, 3.0, 5.5, 3.5)


def output_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Outputs Generated By The System", "What the user gets after clicking Run Battery Intelligence")
    add_panel(slide, 0.7, 1.45, 5.8, 5.6, "Primary Outputs")
    add_bullets(
        slide,
        [
            "State of Health (SoH)",
            "Wear percentage",
            "Remaining Useful Life (RUL)",
            "Estimated end-of-life cycle",
            "Failure probability",
            "Thermal risk label",
            "Confidence score and anomaly flag",
        ],
        0.95,
        1.9,
        5.0,
        4.0,
        font_size=16,
    )
    add_panel(slide, 6.75, 1.45, 5.8, 5.6, "Analytical Visuals")
    add_bullets(
        slide,
        [
            "Capacity degradation forecast chart",
            "Feature contribution / wear-driver chart",
            "Pack-level heatmap",
            "What-if low-stress scenario",
            "Maintenance recommendations",
            "Downloadable Excel report",
        ],
        7.0,
        1.9,
        5.0,
        4.0,
        font_size=16,
    )


def ui_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "UI And User Experience", "How the dashboard helps the user understand battery condition")
    add_bullets(
        slide,
        [
            "Top cards summarize dataset strength and model performance.",
            "Interactive controls let the user simulate EV or drone battery scenarios.",
            "The health ring gives a quick visual signal of battery condition.",
            "Charts explain both future degradation and the main reasons behind wear.",
            "The report export makes the result easy to share and inspect in Excel.",
        ],
        0.9,
        1.7,
        5.5,
        3.7,
        font_size=16,
    )
    add_panel(slide, 6.6, 1.55, 5.8, 4.8, "Why this matters")
    add_bullets(
        slide,
        [
            "Transforms raw ML output into an interpretable battery intelligence tool",
            "Useful for predictive maintenance demos and project presentations",
            "Balances engineering realism with an accessible visual interface",
        ],
        6.9,
        2.0,
        5.0,
        2.8,
        font_size=16,
    )


def limitations_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Limitations And Future Work", "What is realistic already and what can be improved further")
    add_panel(slide, 0.7, 1.45, 5.8, 5.7, "Current Limitations")
    add_bullets(
        slide,
        [
            "NASA dataset is real, but some EV/drone mission variables are not directly present in the raw data.",
            "Several scenario variables are modeled through domain logic rather than learned from large telemetry datasets.",
            "The system is strong for demonstration and analysis, but not yet a production battery management system.",
        ],
        0.95,
        1.95,
        5.0,
        3.8,
        font_size=16,
    )
    add_panel(slide, 6.8, 1.45, 5.6, 5.7, "Future Improvements")
    add_bullets(
        slide,
        [
            "Use larger EV/drone telemetry datasets for direct learning of mission stress.",
            "Add sequence models such as LSTM or transformer-based time-series predictors.",
            "Include real cell-balancing, fault-code, and BMS telemetry streams.",
            "Deploy the app as a cloud dashboard or desktop product.",
        ],
        7.05,
        1.95,
        4.9,
        3.8,
        font_size=16,
    )


def conclusion_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Conclusion", "Why this project is meaningful")
    add_panel(slide, 0.8, 1.7, 11.8, 4.8)
    add_bullets(
        slide,
        [
            "This project combines real NASA battery-aging data with machine learning and an engineering-based realism layer.",
            "It predicts battery degradation, explains the result visually, and supports EV/drone-specific what-if analysis.",
            "The final system is not just a model, but a complete battery intelligence dashboard with reporting and interpretability.",
            "That makes it a strong academic project, a good portfolio piece, and a useful demo for predictive maintenance concepts.",
        ],
        1.1,
        2.15,
        10.9,
        3.2,
        font_size=20,
    )
    note = slide.shapes.add_textbox(Inches(0.85), Inches(6.7), Inches(11.5), Inches(0.35))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Presentation generated from the current project artifacts in C:\\ML proj"
    r.font.size = Pt(12)
    r.font.color.rgb = MUTED


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover_slide(prs)
    abbreviations_slide(prs)
    workflow_slide(prs)
    data_slide(prs)
    parameters_slide(prs)
    formulas_slide(prs)
    algorithms_slide(prs)
    results_slide(prs)
    output_slide(prs)
    ui_slide(prs)
    limitations_slide(prs)
    conclusion_slide(prs)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(path)
